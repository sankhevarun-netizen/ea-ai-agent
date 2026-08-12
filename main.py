import base64
import os
import zlib
from pathlib import Path
from typing import Optional, Dict, Any

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.middleware.gzip import GZipMiddleware
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles

# Load frontend HTML at startup so Vercel Lambda can serve it from memory
_FRONTEND_HTML = ""
try:
    _frontend_path = Path(__file__).parent / "frontend" / "index.html"
    if _frontend_path.exists():
        _FRONTEND_HTML = _frontend_path.read_text(encoding="utf-8")
except Exception:
    pass

_ASSESSMENT_HTML = ""
for _ap in [
    Path(__file__).parent / "frontend" / "ea_assessment.html",
    Path("frontend") / "ea_assessment.html",
    Path("/var/task/frontend/ea_assessment.html"),
]:
    try:
        if _ap.exists():
            _ASSESSMENT_HTML = _ap.read_text(encoding="utf-8")
            break
    except Exception:
        continue

from models.schemas import EARequest, EAResponse
from orchestrator import route_request, resolve_agents
from agents.arb_agent import run_arb
from agents.time_agent import run_time, _score_tools, _extract_tools_from_text, _extract_tools_from_image
from agents.mapping_agent import run_mapping
from agents.maturity_agent import run_maturity
from agents.insights_agent import run_insights
from agents.full_pipeline import run_full_pipeline
from utils.duplication_detector import DuplicationDetector
from utils.report_generator import ReportGenerator

_on_vercel = os.getenv("VERCEL") == "1"
REPORTS_DIR = Path("/tmp/ea-reports") if _on_vercel else Path(__file__).parent / "reports"
REPORTS_DIR.mkdir(parents=True, exist_ok=True)

_detector = DuplicationDetector()
_reporter = ReportGenerator()

app = FastAPI(
    title="EA AI Intelligence",
    description="Multi-agent Enterprise Architecture AI platform — ARB, TIME, Mapping, Maturity, Insights",
    version="1.0.0",
)

app.add_middleware(GZipMiddleware, minimum_size=1000)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

AGENT_MAP = {
    "ARB": run_arb,
    "TIME": run_time,
    "MAPPING": run_mapping,
    "MATURITY": run_maturity,
    "INSIGHTS": run_insights,
}

_FRONTEND_DIR = Path(__file__).parent / "frontend"
if _FRONTEND_DIR.exists() and not _on_vercel:
    app.mount("/ui", StaticFiles(directory=str(_FRONTEND_DIR), html=True), name="frontend")


# ─── Root — serve frontend ────────────────────────────────────────────────────

@app.get("/", response_class=HTMLResponse, include_in_schema=False)
def serve_ui():
    if _FRONTEND_HTML:
        return HTMLResponse(content=_FRONTEND_HTML)
    return HTMLResponse(content="<h2>EA AI Intelligence — <a href='/docs'>See API Docs</a></h2>")


@app.get("/assessment", response_class=HTMLResponse, include_in_schema=False)
def serve_assessment():
    if _ASSESSMENT_HTML:
        return HTMLResponse(content=_ASSESSMENT_HTML)
    return HTMLResponse(content="<h2>EA Maturity Assessment — file not found</h2>")


# ─── Core EA Agent endpoint ───────────────────────────────────────────────────

@app.get("/health")
def health():
    claude_status = "unknown"
    claude_error = None
    model_used = None
    try:
        from services.claude_service import _get_client, MODEL
        client = _get_client()
        client.messages.create(
            model=MODEL,
            max_tokens=1,
            messages=[{"role": "user", "content": "ping"}],
        )
        claude_status = "ok"
        model_used = MODEL
    except Exception as e:
        claude_status = "error"
        claude_error = str(e)

    return {
        "status": "EA AI Intelligence running",
        "version": "1.0.0",
        "agents": list(AGENT_MAP.keys()),
        "claude": claude_status,
        "model": model_used,
        "claude_error": claude_error,
    }


@app.post("/ea-agent", response_model=EAResponse)
def handle_request(request: EARequest):
    try:
        if request.agent_override:
            override = request.agent_override.upper()
            if override not in AGENT_MAP:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unknown agent '{override}'. Valid: {list(AGENT_MAP.keys())}",
                )
            agents_to_run = [override]
        else:
            intent = route_request(request.query)
            agents_to_run = resolve_agents(intent, request.query)

        results: dict = {}
        for agent_name in agents_to_run:
            if agent_name in AGENT_MAP:
                results[agent_name] = AGENT_MAP[agent_name](request.data or {})

        if len(agents_to_run) > 1 and "INSIGHTS" not in agents_to_run:
            results["INSIGHTS"] = run_insights(results)
            agents_to_run = agents_to_run + ["INSIGHTS"]

        return EAResponse(
            agent_used=agents_to_run[0] if len(agents_to_run) == 1 else "MULTI",
            result=results,
            agents_chain=agents_to_run,
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── TIME / Portfolio endpoints ───────────────────────────────────────────────

# ── ARB proposal extraction prompt ───────────────────────────────────────────
_ARB_EXTRACT_PROMPT = """You are an Enterprise Architecture governance expert.
Extract architecture proposal details from the provided document, image or text.
Return ONLY a valid JSON object with exactly these fields (use null for anything not found):
{
  "proposal_title": "...",
  "description": "...",
  "technology_stack": ["tech1", "tech2"],
  "timeline": "...",
  "estimated_cost": 0,
  "annual_running_cost": 0,
  "business_case": "...",
  "current_system": "...",
  "integration_points": "...",
  "security_compliance": "...",
  "data_sensitivity": "...",
  "alternatives_considered": "..."
}
Extract every detail you can find. Do not invent information that is not in the source."""


@app.post("/arb/ingest")
async def arb_ingest(
    file: Optional[UploadFile] = File(None),
    free_text: Optional[str] = Form(None),
):
    """
    Extract structured ARB proposal fields from an uploaded document or image.
    Supports: PDF, PPTX, DOCX, XLSX, XLS, CSV, TXT, MD, PNG, JPG, JPEG, GIF, WEBP.
    Returns JSON with proposal_title, description, technology_stack, costs, etc.
    """
    from services.claude_service import call_claude, call_claude_vision
    import json as _json

    def _parse_result(raw: str) -> dict:
        clean = raw.strip().lstrip("```json").lstrip("```").rstrip("```").strip()
        return _json.loads(clean)

    # ── IMAGE / Vision path ───────────────────────────────────────────────────
    if file and file.filename:
        content = await file.read()
        ext = Path(file.filename).suffix.lower()
        media_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                     ".gif": "image/gif", ".webp": "image/webp"}

        if ext in media_map:
            raw = call_claude_vision(
                _ARB_EXTRACT_PROMPT, content, media_map[ext],
                "Extract all architecture proposal details visible in this image. Return structured JSON."
            )
            try:
                return _parse_result(raw)
            except Exception:
                raise HTTPException(422, f"Could not parse image extraction: {raw[:300]}")

        # ── PDF ───────────────────────────────────────────────────────────────
        if ext == ".pdf":
            import pypdf, io as _io
            reader = pypdf.PdfReader(_io.BytesIO(content))
            extracted = "\n".join(page.extract_text() or "" for page in reader.pages)
            if not extracted.strip():
                # Scanned / image-based PDF — run vision on each page image
                for page in reader.pages:
                    for img_obj in page.images:
                        raw = call_claude_vision(
                            _ARB_EXTRACT_PROMPT, img_obj.data, "image/png",
                            "This is a page from an architecture proposal document. Extract all proposal details."
                        )
                        try:
                            return _parse_result(raw)
                        except Exception:
                            continue
                raise HTTPException(422, "PDF appears to be image-only and no page images could be read.")

        # ── PPTX / PPT ────────────────────────────────────────────────────────
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            import io as _io
            prs = Presentation(_io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            extracted = "\n".join(texts)

        # ── DOCX ─────────────────────────────────────────────────────────────
        elif ext in (".docx", ".doc"):
            try:
                import docx as _docx, io as _io
                doc = _docx.Document(_io.BytesIO(content))
                extracted = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except Exception:
                extracted = content.decode("utf-8", errors="replace")

        # ── Excel / CSV ──────────────────────────────────────────────────────
        elif ext in (".xlsx", ".xls", ".csv"):
            import pandas as pd, io as _io
            try:
                if ext == ".csv":
                    sheets = {"Sheet1": pd.read_csv(_io.BytesIO(content))}
                else:
                    sheets = pd.read_excel(_io.BytesIO(content), sheet_name=None)
                parts = []
                for sheet_name, df in sheets.items():
                    if df.empty:
                        continue
                    parts.append(f"--- Sheet: {sheet_name} ---")
                    parts.append(df.to_csv(index=False))
                extracted = "\n".join(parts)
            except Exception as e:
                raise HTTPException(422, f"Could not read spreadsheet: {e}")

        # ── Plain text / Markdown ─────────────────────────────────────────────
        else:
            extracted = content.decode("utf-8", errors="replace")

    elif free_text:
        extracted = free_text
    else:
        raise HTTPException(400, "No file or text provided.")

    if not extracted.strip():
        raise HTTPException(422, "No readable content found in the uploaded file.")

    raw = call_claude(
        _ARB_EXTRACT_PROMPT,
        f"Extract the ARB proposal details from this document:\n\n{extracted[:8000]}"
    )
    try:
        return _parse_result(raw)
    except Exception:
        raise HTTPException(422, f"Could not structure the extracted content. Raw: {raw[:300]}")


@app.post("/arb/evaluate")
async def arb_evaluate(body: Dict[str, Any]):
    """
    Run a full ARB governance decision using portfolio rationalization results.
    Auto-builds the proposal from flagged applications, then calls the ARB agent.

    Body keys:
      applications          — full list from /time/ingest
      proposal_title        — optional override
      business_case         — optional override
      timeline              — optional override (default "12-18 months")
      ea_context            — context dict from frontend wizard
      rationalization_summary — summary dict from /time/ingest
    """
    from agents.arb_agent import run_arb

    applications = body.get("applications", [])
    flagged = [
        a for a in applications
        if a.get("rationalization_action") in ("Retire", "Replace", "Refactor", "Rehost", "Replatform")
        or a.get("rationalization_action") in ("Replace", "Retire", "Refactor")
    ]

    # Build summary of flagged apps for the proposal
    flagged_names = [a.get("name", "Unknown") for a in flagged]
    vendors = list({a.get("vendor", "") for a in applications if a.get("vendor")})[:8]
    total_cost = sum(a.get("annual_cost") or 0 for a in applications)
    flagged_cost = sum(a.get("annual_cost") or 0 for a in flagged)

    ea_ctx = body.get("ea_context") or {}
    rat_summary = body.get("rationalization_summary") or {}

    proposal = {
        "proposal_title": body.get("proposal_title") or
            f"Portfolio Rationalization — {len(flagged)} Application{'s' if len(flagged)!=1 else ''} Requiring Architectural Action",
        "description": (
            f"EA Rationalization analysis of {len(applications)} applications identifies "
            f"{len(flagged)} requiring architectural change: "
            f"{', '.join(flagged_names[:6])}{'...' if len(flagged_names)>6 else ''}. "
            f"Annual portfolio spend: ${total_cost:,.0f}. "
            f"Flagged application spend: ${flagged_cost:,.0f}."
        ),
        "technology_stack": vendors,
        "timeline": body.get("timeline") or "12-18 months",
        "estimated_cost": body.get("estimated_cost") or 0,
        "annual_running_cost": body.get("annual_running_cost") or 0,
        "business_case": body.get("business_case") or (
            f"Portfolio analysis shows {len(flagged)} of {len(applications)} applications are "
            f"flagged for action (Retire/Replace/Refactor/Rehost/Replatform). Rationalization will reduce technical debt, "
            f"eliminate duplications ({rat_summary.get('duplications_found', 0)} found), "
            f"and align the portfolio with strategic business outcomes. "
            f"Industry: {ea_ctx.get('industry', 'not specified')}. "
            f"EA framework: {ea_ctx.get('ea_framework', 'not specified')}. "
            f"{ea_ctx.get('additional_context', '')}"
        ).strip(),
        "current_system": f"Portfolio of {len(applications)} applications",
        "integration_points": body.get("integration_points") or
            f"Applications span {len(set(a.get('category','') for a in applications if a.get('category')))} categories",
        "security_compliance": ea_ctx.get("governance") or "To be assessed",
        "data_sensitivity": body.get("data_sensitivity") or "Mixed — includes Critical and High business systems",
        "alternatives_considered": body.get("alternatives_considered") or
            "Maintain status quo (rejected — unsustainable cost and technical debt trajectory)",
        "flagged_applications": [
            {
                "name": a.get("name"),
                "vendor": a.get("vendor"),
                "category": a.get("category"),
                "rationalization_action": a.get("rationalization_action"),
                "rationalization_action": a.get("rationalization_action"),
                "annual_cost": a.get("annual_cost"),
                "composite_score": a.get("composite_score"),
                "end_of_life": a.get("end_of_life"),
            }
            for a in flagged
        ],
        "full_portfolio_summary": rat_summary,
        "ea_context": ea_ctx,
    }

    try:
        result = run_arb(proposal)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"ARB evaluation failed: {str(e)}")

    # If run_arb returned a parse error, try to re-extract JSON from raw_response
    if result.get("parse_error") and result.get("raw_response"):
        import re as _re, json as _json2
        raw = result["raw_response"]
        # Try to find JSON block in the raw response
        json_match = _re.search(r'\{[\s\S]*\}', raw)
        if json_match:
            try:
                result = _json2.loads(json_match.group())
            except Exception:
                pass

    # Normalise decision field
    decision = str(result.get("decision", "")).upper()
    if decision not in ("APPROVED", "REJECTED", "CONDITIONAL"):
        decision = "CONDITIONAL"

    return {
        "decision": decision,
        "confidence_score": result.get("confidence_score") or result.get("confidence") or 0,
        "compliance_score": result.get("compliance_score") or result.get("compliance") or 0,
        "risks": result.get("risks") or [],
        "anti_patterns_detected": result.get("anti_patterns_detected") or [],
        "redundancy_flags": result.get("redundancy_flags") or [],
        "integration_complexity": result.get("integration_complexity") or "UNKNOWN",
        "recommendations": result.get("recommendations") or [],
        "conditions": result.get("conditions") or [],
        "reasoning": result.get("reasoning") or "",
        "proposal": proposal,
        "flagged_count": len(flagged),
        "total_apps": len(applications),
    }


@app.post("/arb/signoff")
async def arb_signoff(body: Dict[str, Any]):
    """
    Store the human sign-off audit trail for an ARB decision.
    Called after architect review and/or CIO approval in the pipeline wizard.

    Body keys:
      arb_decision        — original AI ARB result dict
      architect_name      — reviewer name
      architect_decision  — 'accept' | 'override_approved' | 'override_conditional' | 'override_rejected'
      architect_notes     — free-text review comments
      cio_name            — approver name
      cio_decision        — 'approved' | 'rejected' | 'more_info'
      cio_notes           — sign-off notes
      signed_at           — ISO timestamp from client
    """
    from services.supabase_service import store_decision
    from datetime import datetime, timezone

    arb_decision    = body.get("arb_decision", {})
    architect_name  = (body.get("architect_name") or "").strip()
    arch_decision   = body.get("architect_decision", "accept")
    architect_notes = (body.get("architect_notes") or "").strip()
    cio_name        = (body.get("cio_name") or "").strip()
    cio_decision    = body.get("cio_decision", "")
    cio_notes       = (body.get("cio_notes") or "").strip()
    signed_at       = body.get("signed_at") or datetime.now(timezone.utc).isoformat()

    # Resolve effective final decision (CIO may override architect, architect may override AI)
    ai_decision = (arb_decision.get("decision") or "UNKNOWN").upper()
    override_map = {
        "override_approved":    "APPROVED",
        "override_conditional": "CONDITIONAL",
        "override_rejected":    "REJECTED",
    }
    effective_decision = override_map.get(arch_decision, ai_decision)
    if cio_decision == "rejected":
        effective_decision = "REJECTED"
    elif cio_decision == "more_info":
        effective_decision = "PENDING_MORE_INFO"

    signoff = {
        "architect": {
            "name":     architect_name,
            "decision": arch_decision,
            "notes":    architect_notes,
        },
        "cio": {
            "name":     cio_name,
            "decision": cio_decision,
            "notes":    cio_notes,
        } if cio_name else None,
        "effective_decision": effective_decision,
        "ai_decision":        ai_decision,
        "signed_at":          signed_at,
    }

    store_decision(
        {"arb_decision": ai_decision, "architect": architect_name, "cio": cio_name},
        __import__("json").dumps(signoff),
        agent="ARB_SIGNOFF",
    )

    return {
        **arb_decision,
        "signoff": signoff,
        "effective_decision": effective_decision,
    }


@app.post("/time/ingest")
async def time_ingest(
    file: Optional[UploadFile] = File(None),
    free_text: Optional[str] = Form(None),
    applications: Optional[str] = Form(None),  # JSON string of tool list
    industry: Optional[str] = Form(None),
    sub_sector: Optional[str] = Form(None),
):
    """
    Ingest portfolio data (file upload, free text, or JSON body),
    score every tool with the ScoringEngine, detect duplications,
    and return the enriched portfolio ready for report generation.
    """
    tools_raw = []

    if file and file.filename:
        content = await file.read()
        ext = Path(file.filename).suffix.lower()

        if ext == ".json":
            import json
            data = json.loads(content)
            tools_raw = data if isinstance(data, list) else list(data.values())[0] if data else []

        elif ext in (".csv",):
            import csv, io
            # Try utf-8 first, fall back to latin-1 for Excel-exported CSVs with £/€ symbols
            for enc in ("utf-8-sig", "utf-8", "latin-1"):
                try:
                    decoded = content.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                decoded = content.decode("utf-8", errors="replace")
            reader = csv.DictReader(io.StringIO(decoded))
            tools_raw = [
                {k.strip(): (v.strip() if isinstance(v, str) else v)
                 for k, v in row.items() if k}
                for row in reader
                if any(v and str(v).strip() not in ("", "nan", "None") for v in row.values())
            ]

        elif ext in (".xlsx", ".xls"):
            import tempfile, pandas as pd, json as _json
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            # ── Auto-detect header row ────────────────────────────────────────
            # Some templates have a title/subtitle row before the real headers.
            # We check the first 5 rows: the header row is the one that contains
            # recognisable EA column keywords.
            _EA_KEYWORDS = {
                "name","application","app","tool","system","software","product",
                "vendor","supplier","category","type","cost","spend","budget",
                "users","user","deployment","hosting","age","criticality","priority",
                "compliance","eol","end of life","integrations","business unit","department",
            }
            _header_row = 0
            try:
                for _try_row in range(5):
                    _df_try = pd.read_excel(tmp_path, header=_try_row, nrows=1)
                    # Only consider real columns — skip pandas "Unnamed: N" placeholders
                    _real_cols = [
                        str(c).strip().lower() for c in _df_try.columns
                        if not str(c).strip().lower().startswith("unnamed:")
                    ]
                    # Need at least 3 named columns for this to be a real header row
                    if len(_real_cols) < 3:
                        continue
                    _hits = sum(1 for k in _EA_KEYWORDS
                                if any(c == k or c.startswith(k + " ") or c.endswith(" " + k) or (" " + k + " ") in c
                                       for c in _real_cols))
                    if _hits >= 3:
                        _header_row = _try_row
                        break
            except Exception:
                _header_row = 0
            df = pd.read_excel(tmp_path, header=_header_row)
            os.unlink(tmp_path)
            # Replace all float NaN with None before JSON serialisation
            df = df.where(pd.notnull(df), None)
            tools_raw = _json.loads(df.to_json(orient="records"))
            # Filter out entirely-blank rows
            tools_raw = [r for r in tools_raw if any(
                v is not None and str(v).strip() not in ("", "nan", "None")
                for v in r.values()
            )]

        elif ext == ".pdf":
            import pypdf, io as _io
            reader = pypdf.PdfReader(_io.BytesIO(content))
            extracted = "\n".join(page.extract_text() or "" for page in reader.pages)
            if extracted.strip():
                tools_raw = _extract_tools_from_text(extracted)
            else:
                # Scanned/image-only PDF — use vision on each page image
                page_tools: list = []
                for page in reader.pages:
                    for img_obj in page.images:
                        page_tools.extend(_extract_tools_from_image(img_obj.data, "image/png"))
                tools_raw = page_tools

        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            import io as _io
            prs = Presentation(_io.BytesIO(content))
            slide_texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
            tools_raw = _extract_tools_from_text("\n".join(slide_texts))

        elif ext in (".docx", ".doc"):
            try:
                import docx as _docx, io as _io
                doc = _docx.Document(_io.BytesIO(content))
                doc_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                # Also extract tables (often contain app inventories)
                for table in doc.tables:
                    headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
                    for row in table.rows[1:]:
                        row_data = {headers[i]: cell.text.strip()
                                    for i, cell in enumerate(row.cells)
                                    if i < len(headers) and headers[i]}
                        if any(row_data.values()):
                            tools_raw.append(row_data)
                if not tools_raw:
                    tools_raw = _extract_tools_from_text(doc_text)
            except Exception:
                tools_raw = _extract_tools_from_text(content.decode("utf-8", errors="replace"))

        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
            media_map = {
                ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".gif": "image/gif", ".webp": "image/webp",
            }
            tools_raw = _extract_tools_from_image(content, media_map[ext])

        elif ext in (".txt", ".md"):
            tools_raw = _extract_tools_from_text(content.decode("utf-8", errors="replace"))

        else:
            # Generic fallback — try text extraction
            tools_raw = _extract_tools_from_text(content.decode("utf-8", errors="replace"))

    elif free_text:
        tools_raw = _extract_tools_from_text(free_text)

    elif applications:
        import json
        tools_raw = json.loads(applications)

    # ── EA data validation ────────────────────────────────────────────────────
    if not tools_raw:
        raise HTTPException(
            400,
            "No application data found in the uploaded file. "
            "Please upload a file containing your application inventory "
            "(CSV, Excel, PDF, PPTX, image, or DOCX with application names and details)."
        )

    # Check that at least some rows have any non-blank content
    # (tools_raw still has original column names — normalization happens in _score_tools)
    _NAME_KEYS = {"name","application name","application","app name","app","tool name",
                  "tool","system name","system","software","product","product name",
                  "service","solution","platform"}
    def _has_name(row: dict) -> bool:
        for k, v in row.items():
            kl = str(k).strip().lower()
            # Accept exact match OR any name key appearing as substring of column
            # (handles templates that append " *", " (required)", etc.)
            matched = kl in _NAME_KEYS or any(nk in kl for nk in _NAME_KEYS)
            if matched and v and str(v).strip().lower() not in ("", "nan", "none", "unknown"):
                return True
        return False

    named = [t for t in tools_raw if _has_name(t)]
    if not named:
        raise HTTPException(
            422,
            "The file was read successfully but no recognisable application or tool names were found. "
            "Ensure your file contains at least an 'Application Name' or 'Name' column. "
            "Supported formats: CSV, Excel (xlsx/xls), PDF, PPTX, DOCX, PNG, JPG, WEBP."
        )

    scored = _score_tools(tools_raw)
    duplications = _detector.detect_duplications(scored)

    return {
        "success": True,
        "tools_scored": len(scored),
        "industry": industry or "",
        "sub_sector": sub_sector or "",
        "applications": scored,
        "duplications": duplications,
        "summary": {
            "total_apps": len(scored),
            "total_annual_cost": sum(t.get("annual_cost") or 0 for t in scored),
            "duplications_found": len(duplications),
            "potential_savings": sum(d.get("potential_annual_savings", 0) for d in duplications),
            "action_breakdown": _action_counts(scored),
        },
    }


@app.post("/time/report")
async def time_report(body: Dict[str, Any]):
    """
    Generate a PDF portfolio rationalization report.

    Body:
      applications  — list of scored tool dicts (from /time/ingest)
      duplications  — list of duplication records  (from /time/ingest)
      assessments   — optional list of Claude narrative assessments
    """
    tools = body.get("applications", [])
    duplications = body.get("duplications", [])
    assessments = body.get("assessments", [])

    if not tools:
        raise HTTPException(400, "applications list is required.")

    report_path = _reporter.generate(
        tools=tools,
        duplications=duplications,
        assessments=assessments,
        output_dir=str(REPORTS_DIR),
        fmt="pdf",
    )

    return FileResponse(
        path=report_path,
        filename=Path(report_path).name,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{Path(report_path).name}"'},
    )


@app.post("/time/full")
async def time_full(body: Dict[str, Any]):
    """
    One-shot: score portfolio + Claude TIME narrative. Returns JSON only (no PDF).
    Body: { applications: [...] } or { free_text: "..." }
    """
    result = run_time(body)
    return result


# ─── Full EA Pipeline ─────────────────────────────────────────────────────────

@app.post("/ea-full")
async def ea_full(body: Dict[str, Any]):
    """
    Full 4-stage EA Intelligence pipeline:
      TIME → MAPPING → MATURITY → INSIGHTS → PDF report

    Body: { applications: [...] }  or  { free_text: "..." }

    Returns:
      - pipeline JSON with all 4 agent outputs
      - report_url: path to the downloadable PDF report
    """
    try:
        pipeline = run_full_pipeline(body)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Pipeline execution failed: {str(e)}")

    try:
        report_path = _reporter.generate_pipeline_pdf(
            pipeline=pipeline,
            output_dir=str(REPORTS_DIR),
        )
        pipeline["report_path"] = report_path
        pipeline["report_filename"] = Path(report_path).name
        with open(report_path, "rb") as f:
            pdf_bytes = f.read()
        compressed = zlib.compress(pdf_bytes, level=9)
        pipeline["report_b64"] = base64.b64encode(compressed).decode("ascii")
        pipeline["report_compressed"] = True
        pipeline["report_size_kb"] = round(len(pdf_bytes) / 1024, 1)
    except Exception as e:
        pipeline["report_error"] = str(e)
        pipeline["report_filename"] = None
        pipeline["report_b64"] = None

    try:
        html_str = _reporter.generate_pipeline_html(pipeline)
        pipeline["html_b64"] = base64.b64encode(html_str.encode("utf-8")).decode("ascii")
    except Exception as e:
        pipeline["html_b64"] = None

    return pipeline


@app.post("/mapping/analyse")
async def mapping_analyse(body: Dict[str, Any]):
    """
    Run ONLY the MAPPING stage (dependency analysis) for the given applications.
    Called from the pipeline stepper Step 3 so Steps 4-6 can run separately.

    Body keys:
      applications   — full scored list from /time/ingest
      arb_result     — ARB decision output (optional)
      ea_context     — context dict
      mapping_context — mapping-specific context
    """
    from agents.mapping_agent import run_mapping_batch
    import json as _json

    applications = body.get("applications", [])
    arb_result   = body.get("arb_result") or {}
    ea_ctx       = body.get("ea_context") or {}
    map_ctx      = body.get("mapping_context") or {}

    flagged = [
        a for a in applications
        if a.get("rationalization_action") in ("Retire", "Replace", "Refactor", "Rehost", "Replatform")
        or a.get("rationalization_action") in ("Replace", "Retire", "Refactor")
    ]

    if not flagged:
        return {
            "skipped": True,
            "reason": "No applications flagged for action (Retire/Replace/Refactor) — no dependency risk to analyse.",
            "overall_impact_level": "LOW",
            "mapped_applications": [],
            "total_dependencies": 0,
        }

    # Build ARB context string
    arb_summary = ""
    if arb_result:
        arb_summary = (
            f"ARB Decision: {arb_result.get('decision','NOT_RUN')}. "
            f"Compliance: {arb_result.get('compliance_score',0)}. "
            f"Key risks: {'; '.join((arb_result.get('risks') or [])[:3])}."
        )

    additional_ctx = ea_ctx.get("additional_context", "")
    map_notes = map_ctx.get("additional_notes", "")
    combined_ctx = "\n".join(filter(None, [additional_ctx, arb_summary, map_notes]))

    mapping_input = {
        "applications":           flagged,
        "all_apps":               applications,
        "industry":               ea_ctx.get("industry", ""),
        "sub_sector":             ea_ctx.get("sub_sector", ""),
        "governance":             ea_ctx.get("governance", ""),
        "ea_framework":           ea_ctx.get("ea_framework", ""),
        "cloud_target":           ea_ctx.get("cloud_strategy", "") or map_ctx.get("cloud_target", ""),
        "known_integrations":     map_ctx.get("known_integrations", ""),
        "planned_changes":        map_ctx.get("planned_changes", ""),
        "data_sensitivity":       map_ctx.get("data_sensitivity", ""),
        "integration_complexity": map_ctx.get("integration_complexity", ""),
        "cmdb_tool":              map_ctx.get("cmdb_tool", ""),
        "additional_context":     combined_ctx,
    }

    try:
        result = run_mapping_batch(mapping_input)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Mapping analysis failed: {str(e)}")

    return result


@app.post("/compliance/check")
async def compliance_check(body: Dict[str, Any]):
    """
    Run the Compliance Mapping agent against the portfolio for selected regulations.
    This is Phase 2 of Step 3 (Mapping) in the pipeline wizard.

    Body keys:
      applications  — list of scored apps (from TIME output)
      regulations   — list of regulation strings the user toggled ON
      mapping_output — result from /mapping/analyse (or null if skipped)
      industry, sub_sector, governance, additional_context
    """
    from agents.compliance_agent import run_compliance
    try:
        result = run_compliance(body)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Compliance check failed: {str(e)}")
    return result


@app.post("/ea-pipeline/step")
async def pipeline_wizard_step(body: Dict[str, Any]):
    """
    Step-by-step wizard endpoint: accepts pre-computed TIME + MAPPING results,
    runs MATURITY → INSIGHTS, generates PDF.  Called from the frontend wizard
    after Steps 1 (TIME) and 2 (MAPPING) are already shown to the user.

    Body keys:
      time_result        — output from /time/ingest (full JSON)
      mapping_result     — output from MAPPING agent (or null)
      compliance_result  — output from /compliance/check (or null if skipped)
      assessment_results — questionnaire answers from inline questionnaire
      industry, sub_sector, governance, ea_framework, ea_tools,
      cloud_strategy, additional_context
    """
    from agents.full_pipeline import run_pipeline_from_step3
    try:
        pipeline = run_pipeline_from_step3(body)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Wizard pipeline failed: {str(e)}")

    try:
        report_path = _reporter.generate_pipeline_pdf(
            pipeline=pipeline,
            output_dir=str(REPORTS_DIR),
        )
        pipeline["report_filename"] = Path(report_path).name
        with open(report_path, "rb") as f:
            pdf_bytes = f.read()
        compressed = zlib.compress(pdf_bytes, level=9)
        pipeline["report_b64"] = base64.b64encode(compressed).decode("ascii")
        pipeline["report_compressed"] = True
        pipeline["report_size_kb"] = round(len(pdf_bytes) / 1024, 1)
    except Exception as e:
        pipeline["report_error"] = str(e)
        pipeline["report_filename"] = None
        pipeline["report_b64"] = None

    try:
        html_str = _reporter.generate_pipeline_html(pipeline)
        pipeline["html_b64"] = base64.b64encode(html_str.encode("utf-8")).decode("ascii")
    except Exception as e:
        pipeline["html_b64"] = None

    return pipeline


@app.get("/ea-full/report/{filename}")
async def download_pipeline_report(filename: str):
    """Download a previously generated full-pipeline PDF report."""
    report_path = REPORTS_DIR / filename
    if not report_path.exists():
        raise HTTPException(404, f"Report '{filename}' not found.")
    return FileResponse(
        path=str(report_path),
        filename=filename,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ─── AI Assistant Chat ────────────────────────────────────────────────────────

_EA_CHAT_PROMPT = """You are the EA Intelligence Assistant — a step-by-step guide embedded in the EA AI Intelligence platform. You work like GPS turn-by-turn navigation: always tell the user EXACTLY what to do next — which page to go to, which button to click, what to type or upload.

GUIDANCE STYLE RULES (mandatory):
- Always respond with a numbered current step, then a clear "NEXT: [exact action]" instruction.
- Never give generic advice. Always name the exact page, button label, or field.
- After each user message, tell them what the NEXT concrete action is.
- If they say "next", "done", "what now", or similar, immediately give the next step.
- Keep responses short (3-6 lines max). No long paragraphs.

STEP-BY-STEP PIPELINE FLOW (6 Steps):
STEP 1 — Upload application inventory (CSV/Excel/PDF) on the EA Analysis page for Portfolio Rationalization.
STEP 2 — ARB Governance: submit for Architecture Review Board approval (or skip to proceed).
STEP 3 — Run Mapping: dependency and impact analysis for flagged applications.
STEP 4 — Complete the 30-question EA Maturity questionnaire, then run the Maturity Assessment.
STEP 5 — View Strategic Insights: AI-synthesised recommendations across all previous steps.
STEP 6 — Preview and Download the full EA Intelligence PDF report.

PLATFORM AGENTS:
- ARB: Architecture Review Board — reviews proposals, returns APPROVE/REJECT/CONDITIONAL + compliance score + risks.
- RATIONALIZATION: Portfolio analysis — scores apps 0-10 on 7 dimensions, assigns 6R action (Retain/Rehost/Replatform/Refactor/Replace/Retire), identifies duplications and savings.
- MAPPING: Dependency Analysis — maps upstream/downstream dependencies, coupling scores, migration sequencing risk.
- MATURITY: EA Maturity Assessment — scores EA practice 1-5 across 8 dimensions (KPMG 4-pillar framework: Vision & Strategy, People, Processes, Technology). Produces High/Medium/Low band per dimension.
- INSIGHTS: Executive Synthesis — C-suite summary, financial impact, KPIs, quick wins.

CSV format: name, vendor, category, annual_cost, user_count, criticality, deployment, integrations, age_years, end_of_life, compliance_required, business_unit
Valid criticality: Critical, High, Medium, Low
Valid 6R actions: Retain, Rehost, Replatform, Refactor, Replace, Retire

Always end your reply with: "NEXT: [exact action the user should take now]" """


@app.post("/ea-report/pdf")
async def generate_report_only(body: Dict[str, Any]):
    """
    Generate PDF + HTML from pre-computed pipeline data — NO agent re-runs.
    Called from the frontend when pipelineReportB64 is missing after a stepper run.
    Body: the pipeline dict (TIME, ARB, MAPPING, MATURITY, INSIGHTS keys).
    """
    pipeline = body
    result: Dict[str, Any] = {}

    try:
        report_path = _reporter.generate_pipeline_pdf(
            pipeline=pipeline,
            output_dir=str(REPORTS_DIR),
        )
        result["report_filename"] = Path(report_path).name
        with open(report_path, "rb") as f:
            pdf_bytes = f.read()
        compressed = zlib.compress(pdf_bytes, level=9)
        result["report_b64"] = base64.b64encode(compressed).decode("ascii")
        result["report_compressed"] = True
        result["report_size_kb"] = round(len(pdf_bytes) / 1024, 1)
    except Exception as e:
        result["report_error"] = str(e)
        result["report_b64"] = None

    try:
        html_str = _reporter.generate_pipeline_html(pipeline)
        result["html_b64"] = base64.b64encode(html_str.encode("utf-8")).decode("ascii")
    except Exception as e:
        result["html_b64"] = None

    return result


@app.post("/chat")
async def chat_assistant(body: Dict[str, Any]):
    """EA Intelligence floating assistant chatbot."""
    from services.claude_service import call_claude_with_history
    messages = body.get("messages", [])
    if not messages:
        raise HTTPException(400, "messages required")
    response = call_claude_with_history(_EA_CHAT_PROMPT, messages)
    return {"response": response}


# ─── Helper ──────────────────────────────────────────────────────────────────

def _action_counts(tools: list) -> dict:
    counts: dict = {}
    for t in tools:
        action = t.get("rationalization_action", "Replatform")
        counts[action] = counts.get(action, 0) + 1
    return counts
